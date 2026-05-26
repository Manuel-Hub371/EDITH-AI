import logging
from pathlib import Path
from typing import Dict, Any

from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)

class DocumentReader:
    """
    Handles extracting text and content from various document formats.
    Supports txt, md, pdf, docx, xlsx, pptx.
    """

    def read_file(self, file_path: str) -> Dict[str, Any]:
        """
        Intelligently routes reading based on file extension.
        """
        path = Path(file_path).resolve()
        
        if not path.exists() or not path.is_file():
            return {"success": False, "error": f"File not found: {path}"}
            
        ext = path.suffix.lower()
        
        try:
            if ext in [".txt", ".md", ".json", ".csv", ".py", ".js", ".html"]:
                content = self.read_text(path)
            elif ext == ".pdf":
                content = self.read_pdf(path)
            elif ext == ".docx":
                content = self.read_docx(path)
            elif ext == ".xlsx":
                content = self.read_xlsx(path)
            elif ext == ".pptx":
                content = self.read_pptx(path)
            else:
                return {"success": False, "error": f"Unsupported file type for reading: {ext}"}
                
            return {
                "success": True,
                "file_path": str(path),
                "extension": ext,
                "content": content
            }
        except Exception as e:
            logger.error(f"Error reading file {path}: {e}")
            return {"success": False, "error": str(e)}

    def read_text(self, path: Path) -> str:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    def read_pdf(self, path: Path) -> str:
        reader = PdfReader(str(path))
        text = ""
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        return text

    def read_docx(self, path: Path) -> str:
        doc = Document(str(path))
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def read_xlsx(self, path: Path) -> str:
        wb = load_workbook(str(path), data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            lines.append(f"--- Sheet: {sheet_name} ---")
            for row in sheet.iter_rows(values_only=True):
                # Filter out completely empty rows
                if any(cell is not None for cell in row):
                    row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    lines.append(row_str)
        return "\n".join(lines)

    def read_pptx(self, path: Path) -> str:
        prs = Presentation(str(path))
        lines = []
        for idx, slide in enumerate(prs.slides):
            lines.append(f"--- Slide {idx + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)

# Singleton
document_reader = DocumentReader()
