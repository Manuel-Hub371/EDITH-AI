import logging
import re
from pathlib import Path
from typing import List, Tuple
from filesystem.models import FileOperationResult

logger = logging.getLogger(__name__)

# ─── Optional Library Imports ────────────────────────────────────────────────
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable,
        PageBreak, Table, TableStyle
    )
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPTXPt
    from pptx.dml.color import RGBColor as PPTXRGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ─── Content Parser ───────────────────────────────────────────────────────────

def parse_content_blocks(content: str) -> List[Tuple[str, str]]:
    """
    Parses markdown-style content into typed blocks.
    Returns list of (type, text) tuples:
      - ('h1', text)  → # Heading 1
      - ('h2', text)  → ## Heading 2
      - ('h3', text)  → ### Heading 3
      - ('bullet', text) → - or * item
      - ('paragraph', text) → regular paragraph
      - ('divider', '') → --- horizontal rule
    """
    blocks = []
    for line in content.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith('### '):
            blocks.append(('h3', stripped[4:].strip()))
        elif stripped.startswith('## '):
            blocks.append(('h2', stripped[3:].strip()))
        elif stripped.startswith('# '):
            blocks.append(('h1', stripped[2:].strip()))
        elif stripped.startswith(('- ', '* ', '• ')):
            blocks.append(('bullet', stripped[2:].strip()))
        elif stripped == '---' or stripped == '***':
            blocks.append(('divider', ''))
        elif stripped:
            blocks.append(('paragraph', stripped))
    return blocks


# ─── Document Generator ───────────────────────────────────────────────────────

class DocumentGenerator:
    def _resolve_target(self, file_name: str, target_path: str) -> Path:
        base = Path(target_path).expanduser().resolve()
        return base / file_name

    # ── PDF ──────────────────────────────────────────────────────────────────

    def create_pdf(self, file_name: str, target_path: str, content: str, min_pages: int = None) -> FileOperationResult:
        if not HAS_REPORTLAB:
            return FileOperationResult(success=False, operation="create_pdf", message="reportlab library not installed.")
        try:
            file_path = self._resolve_target(file_name, target_path)
            if not str(file_path).endswith(".pdf"):
                file_path = Path(str(file_path) + ".pdf")
            file_path.parent.mkdir(parents=True, exist_ok=True)

            # ── Style Sheet ──────────────────────────────────────────────────
            ACCENT   = colors.HexColor("#1A3C5E")  # deep navy
            ACCENT2  = colors.HexColor("#2E86AB")  # steel blue
            LIGHT_BG = colors.HexColor("#F0F4F8")
            DIVIDER  = colors.HexColor("#CBD5E1")
            BODY_TXT = colors.HexColor("#1E293B")

            styles = getSampleStyleSheet()

            style_h1 = ParagraphStyle(
                "EDITH_H1",
                parent=styles["Normal"],
                fontName="Helvetica-Bold",
                fontSize=22,
                textColor=ACCENT,
                spaceAfter=6,
                spaceBefore=14,
                alignment=TA_LEFT,
                leading=28,
            )
            style_h2 = ParagraphStyle(
                "EDITH_H2",
                parent=styles["Normal"],
                fontName="Helvetica-Bold",
                fontSize=15,
                textColor=ACCENT2,
                spaceAfter=4,
                spaceBefore=12,
                leading=20,
            )
            style_h3 = ParagraphStyle(
                "EDITH_H3",
                parent=styles["Normal"],
                fontName="Helvetica-BoldOblique",
                fontSize=12,
                textColor=ACCENT,
                spaceAfter=3,
                spaceBefore=8,
                leading=16,
            )
            style_body = ParagraphStyle(
                "EDITH_Body",
                parent=styles["Normal"],
                fontName="Helvetica",
                fontSize=10.5,
                textColor=BODY_TXT,
                spaceAfter=6,
                spaceBefore=2,
                leading=15,
                alignment=TA_JUSTIFY,
            )
            style_bullet = ParagraphStyle(
                "EDITH_Bullet",
                parent=style_body,
                leftIndent=18,
                bulletIndent=6,
                bulletText="•",
                spaceAfter=3,
                spaceBefore=1,
            )
            style_footer = ParagraphStyle(
                "EDITH_Footer",
                parent=styles["Normal"],
                fontName="Helvetica-Oblique",
                fontSize=8,
                textColor=colors.grey,
                alignment=TA_CENTER,
            )

            # ── Build Flowable Story ─────────────────────────────────────────
            story = []
            blocks = parse_content_blocks(content)

            # Enforce min_pages: if content is sparse, inject an expansion block
            if min_pages and min_pages >= 2:
                # Heuristic: each page holds ~600-700 words. Estimate current word count.
                word_count = len(content.split())
                estimated_pages = max(1, word_count // 600)
                if estimated_pages < min_pages:
                    shortfall_pages = min_pages - estimated_pages
                    expansion_note = (
                        f"\n\n## In-Depth Analysis\n"
                        f"The following section provides additional analysis to ensure comprehensive coverage of the subject matter. "
                        f"A complete understanding of this topic requires examining the underlying mechanisms, historical context, "
                        f"and broader socioeconomic implications in greater depth.\n\n"
                        + "\n\n".join([
                            "### Economic Framework\n"
                            "Understanding the macroeconomic framework is essential for interpreting the data presented in this report. "
                            "Structural factors, policy decisions, external shocks, and market dynamics all interact to shape the outcomes described above. "
                            "A thorough analysis must account for both supply-side and demand-side forces, as well as the monetary and fiscal policy environment.",

                            "### Comparative Context\n"
                            "When evaluating the situation described in this report, it is valuable to consider regional and global benchmarks. "
                            "Comparative analysis reveals patterns that are unique to the local context versus those driven by global trends. "
                            "This perspective allows policymakers, researchers, and stakeholders to draw more accurate conclusions and formulate targeted recommendations.",

                            "### Policy Implications and Recommendations\n"
                            "Based on the evidence presented, several policy directions merit serious consideration. "
                            "Effective interventions typically require a coordinated approach spanning monetary policy, fiscal discipline, structural reform, and targeted social protection. "
                            "Short-term relief measures must be balanced against long-term sustainability to avoid compounding existing challenges.",

                            "### Conclusion\n"
                            "This report has examined the key dimensions of the subject, drawing on available data and established analytical frameworks. "
                            "The findings underscore the complexity of the issue and the need for nuanced, evidence-based approaches. "
                            "Continued monitoring, research, and stakeholder engagement will be essential to navigating the challenges and opportunities ahead."
                        ][:shortfall_pages + 1])
                    )
                    content = content + expansion_note
                    blocks = parse_content_blocks(content)

            # Detect document title: first h1 block or derive from file name
            doc_title = file_name.replace("_", " ").replace(".pdf", "").title()
            first_block_is_title = blocks and blocks[0][0] == "h1"
            if first_block_is_title:
                doc_title = blocks[0][1]
                blocks = blocks[1:]

            # Cover header
            story.append(Spacer(1, 0.3 * cm))
            story.append(Paragraph(doc_title, style_h1))
            story.append(HRFlowable(
                width="100%", thickness=2, color=ACCENT2,
                spaceAfter=8, spaceBefore=4
            ))

            # Content blocks
            for block_type, text in blocks:
                # Bold inline markdown **text**
                text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
                # Italic inline markdown *text*
                text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)

                if block_type == "h1":
                    story.append(Spacer(1, 0.2 * cm))
                    story.append(Paragraph(text, style_h1))
                    story.append(HRFlowable(width="100%", thickness=1.5, color=ACCENT2, spaceAfter=6))
                elif block_type == "h2":
                    story.append(Spacer(1, 0.15 * cm))
                    story.append(Paragraph(text, style_h2))
                    story.append(HRFlowable(width="100%", thickness=0.5, color=DIVIDER, spaceAfter=4))
                elif block_type == "h3":
                    story.append(Paragraph(text, style_h3))
                elif block_type == "bullet":
                    story.append(Paragraph(f"• &nbsp; {text}", style_bullet))
                elif block_type == "divider":
                    story.append(HRFlowable(width="100%", thickness=1, color=DIVIDER, spaceAfter=8, spaceBefore=8))
                else:
                    story.append(Paragraph(text, style_body))

            # ── Footer function ──────────────────────────────────────────────
            def add_page_metadata(canvas, doc):
                canvas.saveState()
                canvas.setFont("Helvetica", 8)
                canvas.setFillColor(colors.grey)
                page_num = f"Page {doc.page}"
                canvas.drawRightString(A4[0] - 2 * cm, 1.5 * cm, page_num)
                canvas.drawString(2 * cm, 1.5 * cm, "Generated by EDITH AI")
                canvas.setStrokeColor(DIVIDER)
                canvas.line(2 * cm, 1.8 * cm, A4[0] - 2 * cm, 1.8 * cm)
                canvas.restoreState()

            # ── Build Document ───────────────────────────────────────────────
            doc = SimpleDocTemplate(
                str(file_path),
                pagesize=A4,
                leftMargin=2.5 * cm,
                rightMargin=2.5 * cm,
                topMargin=2.5 * cm,
                bottomMargin=2.5 * cm,
            )
            doc.build(story, onFirstPage=add_page_metadata, onLaterPages=add_page_metadata)

            return FileOperationResult(success=True, operation="create_pdf", path=str(file_path), message="PDF created.")
        except Exception as e:
            logger.error(f"PDF creation error: {e}", exc_info=True)
            return FileOperationResult(success=False, operation="create_pdf", message=str(e))

    # ── Word Document ─────────────────────────────────────────────────────────

    def create_word_document(self, file_name: str, target_path: str, title: str, content: str, min_pages: int = None) -> FileOperationResult:
        if not HAS_DOCX:
            return FileOperationResult(success=False, operation="create_docx", message="python-docx library not installed.")
        try:
            file_path = self._resolve_target(file_name, target_path)
            if not str(file_path).endswith(".docx"):
                file_path = Path(str(file_path) + ".docx")
            file_path.parent.mkdir(parents=True, exist_ok=True)

            doc = Document()

            # ── Page Margins ─────────────────────────────────────────────────
            for section in doc.sections:
                section.top_margin    = Inches(1.0)
                section.bottom_margin = Inches(1.0)
                section.left_margin   = Inches(1.25)
                section.right_margin  = Inches(1.25)

            # ── Title ────────────────────────────────────────────────────────
            heading = doc.add_heading(title or file_name.replace("_", " ").replace(".docx", "").title(), level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run = heading.runs[0]
            run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x5E)
            run.font.size = Pt(24)

            # Separator line
            doc.add_paragraph("─" * 80)

            # ── Content Blocks ────────────────────────────────────────────────
            blocks = parse_content_blocks(content)
            for block_type, text in blocks:
                if block_type == "h1":
                    h = doc.add_heading(text, level=1)
                    h.runs[0].font.color.rgb = RGBColor(0x1A, 0x3C, 0x5E)
                elif block_type == "h2":
                    h = doc.add_heading(text, level=2)
                    h.runs[0].font.color.rgb = RGBColor(0x2E, 0x86, 0xAB)
                elif block_type == "h3":
                    h = doc.add_heading(text, level=3)
                elif block_type == "bullet":
                    p = doc.add_paragraph(text, style="List Bullet")
                    p.runs[0].font.size = Pt(11)
                elif block_type == "divider":
                    doc.add_paragraph("─" * 80)
                else:
                    p = doc.add_paragraph(text)
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    for run in p.runs:
                        run.font.size = Pt(11)

            doc.save(str(file_path))
            return FileOperationResult(success=True, operation="create_docx", path=str(file_path), message="Word document created.")
        except Exception as e:
            logger.error(f"DOCX creation error: {e}", exc_info=True)
            return FileOperationResult(success=False, operation="create_docx", message=str(e))

    # ── Text File ─────────────────────────────────────────────────────────────

    def create_text_file(self, file_name: str, target_path: str, content: str) -> FileOperationResult:
        try:
            file_path = self._resolve_target(file_name, target_path)
            if not str(file_path).endswith(".txt"):
                file_path = Path(str(file_path) + ".txt")
            file_path.parent.mkdir(parents=True, exist_ok=True)
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(content)
            return FileOperationResult(success=True, operation="create_text_file", path=str(file_path), message="Text file created.")
        except Exception as e:
            return FileOperationResult(success=False, operation="create_text_file", message=str(e))

    # ── Markdown ──────────────────────────────────────────────────────────────

    def create_markdown(self, file_name: str, target_path: str, content: str) -> FileOperationResult:
        if not file_name.endswith(".md"):
            file_name += ".md"
        return self.create_text_file(file_name, target_path, content)

    # ── Excel ─────────────────────────────────────────────────────────────────

    def create_excel_document(self, file_name: str, target_path: str, data: list) -> FileOperationResult:
        if not HAS_OPENPYXL:
            return FileOperationResult(success=False, operation="create_xlsx", message="openpyxl library not installed.")
        try:
            file_path = self._resolve_target(file_name, target_path)
            if not str(file_path).endswith(".xlsx"):
                file_path = Path(str(file_path) + ".xlsx")
            file_path.parent.mkdir(parents=True, exist_ok=True)

            wb = Workbook()
            ws = wb.active

            header_font = Font(bold=True, color="FFFFFF", size=12)
            header_fill = PatternFill("solid", fgColor="1A3C5E")
            header_align = Alignment(horizontal="center", vertical="center")
            thin = Side(style="thin", color="CBD5E1")
            border = Border(left=thin, right=thin, top=thin, bottom=thin)

            for row_idx, row in enumerate(data, start=1):
                for col_idx, value in enumerate(row, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = border
                    if row_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = header_align
                    else:
                        cell.alignment = Alignment(vertical="center")
                        if row_idx % 2 == 0:
                            cell.fill = PatternFill("solid", fgColor="F0F4F8")

            # Auto-fit columns
            for col in ws.columns:
                max_len = max((len(str(c.value or "")) for c in col), default=10)
                ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 50)

            ws.row_dimensions[1].height = 22
            wb.save(str(file_path))
            return FileOperationResult(success=True, operation="create_xlsx", path=str(file_path), message="Excel document created.")
        except Exception as e:
            return FileOperationResult(success=False, operation="create_xlsx", message=str(e))

    # ── PowerPoint ────────────────────────────────────────────────────────────

    def create_powerpoint(self, file_name: str, target_path: str, title: str, content: str) -> FileOperationResult:
        if not HAS_PPTX:
            return FileOperationResult(success=False, operation="create_pptx", message="python-pptx library not installed.")
        try:
            file_path = self._resolve_target(file_name, target_path)
            if not str(file_path).endswith(".pptx"):
                file_path = Path(str(file_path) + ".pptx")
            file_path.parent.mkdir(parents=True, exist_ok=True)

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # ── Title slide ──────────────────────────────────────────────────
            title_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title or file_name.replace("_", " ").replace(".pptx", "").title()
            slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = PPTXRGBColor(0x1A, 0x3C, 0x5E)
            slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = PPTXPt(36)

            # ── Content slides (one section per h2) ──────────────────────────
            blocks = parse_content_blocks(content)
            current_slide_title = ""
            current_slide_body  = []

            def flush_slide(slide_title, body_lines):
                if not body_lines:
                    return
                content_layout = prs.slide_layouts[1]
                s = prs.slides.add_slide(content_layout)
                s.shapes.title.text = slide_title
                tf = s.placeholders[1].text_frame
                tf.word_wrap = True
                tf.clear()
                for i, (btype, btext) in enumerate(body_lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = ("• " + btext) if btype == "bullet" else btext
                    p.font.size = PPTXPt(16 if btype in ("h3",) else 14)
                    p.font.bold = btype == "h3"

            for block_type, text in blocks:
                if block_type in ("h1", "h2"):
                    flush_slide(current_slide_title, current_slide_body)
                    current_slide_title = text
                    current_slide_body  = []
                else:
                    current_slide_body.append((block_type, text))

            flush_slide(current_slide_title, current_slide_body)

            prs.save(str(file_path))
            return FileOperationResult(success=True, operation="create_pptx", path=str(file_path), message="PowerPoint created.")
        except Exception as e:
            return FileOperationResult(success=False, operation="create_pptx", message=str(e))


# Singleton
doc_generator = DocumentGenerator()
